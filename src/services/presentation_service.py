"""
Serviço de geração de apresentações PowerPoint
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import base64
from io import BytesIO


class PresentationService:
    def __init__(self):
        self.default_width = Inches(10)
        self.default_height = Inches(7.5)
    
    def create_presentation(self, title, slides_data):
        """
        Cria uma apresentação PowerPoint
        
        Args:
            title: Título da apresentação
            slides_data: Lista de dicionários com dados dos slides
                [
                    {
                        'title': 'Título do Slide',
                        'content': ['Ponto 1', 'Ponto 2', ...],
                        'layout': 'title_and_content' | 'title_only' | 'blank'
                    }
                ]
        """
        try:
            prs = Presentation()
            prs.slide_width = self.default_width
            prs.slide_height = self.default_height
            
            # Slide de título
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title_shape.text = title
            subtitle.text = "Criado com AI Content Studio"
            
            # Adiciona slides de conteúdo
            for slide_data in slides_data:
                self._add_content_slide(prs, slide_data)
            
            # Salva em memória
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            
            # Converte para base64
            pptx_data = base64.b64encode(output.getvalue()).decode()
            
            return {
                'success': True,
                'data': pptx_data,
                'filename': f'{title.replace(" ", "_")}.pptx',
                'slides_count': len(prs.slides)
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def _add_content_slide(self, prs, slide_data):
        """
        Adiciona um slide de conteúdo à apresentação
        """
        layout_type = slide_data.get('layout', 'title_and_content')
        
        if layout_type == 'title_only':
            slide_layout = prs.slide_layouts[5]  # Title only
        elif layout_type == 'blank':
            slide_layout = prs.slide_layouts[6]  # Blank
        else:
            slide_layout = prs.slide_layouts[1]  # Title and Content
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Define título
        if 'title' in slide_data and slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
            # Estiliza título
            title_frame = slide.shapes.title.text_frame
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Adiciona conteúdo
        if 'content' in slide_data and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            content = slide_data['content']
            if isinstance(content, list):
                for i, item in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(18)
            else:
                p = text_frame.paragraphs[0]
                p.text = content
                p.font.size = Pt(18)
    
    def generate_from_topic(self, topic, num_slides=5):
        """
        Gera uma apresentação básica a partir de um tópico
        """
        try:
            # Estrutura básica de apresentação
            slides_data = []
            
            # Slide de introdução
            slides_data.append({
                'title': 'Introdução',
                'content': [
                    f'Visão geral sobre {topic}',
                    'Conceitos principais',
                    'Importância e aplicações'
                ]
            })
            
            # Slides de conteúdo
            for i in range(2, num_slides):
                slides_data.append({
                    'title': f'Tópico {i-1}',
                    'content': [
                        f'Aspecto {i-1} de {topic}',
                        'Detalhes e características',
                        'Exemplos práticos',
                        'Considerações importantes'
                    ]
                })
            
            # Slide de conclusão
            slides_data.append({
                'title': 'Conclusão',
                'content': [
                    'Resumo dos pontos principais',
                    'Próximos passos',
                    'Recursos adicionais'
                ]
            })
            
            return self.create_presentation(topic, slides_data)
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def create_from_text(self, title, text_content):
        """
        Cria apresentação a partir de texto livre
        Divide o texto em slides automaticamente
        """
        try:
            # Divide texto em parágrafos
            paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
            
            slides_data = []
            current_slide = None
            
            for para in paragraphs:
                # Se o parágrafo é curto, pode ser um título
                if len(para) < 100 and not para.endswith('.'):
                    if current_slide:
                        slides_data.append(current_slide)
                    current_slide = {
                        'title': para,
                        'content': []
                    }
                else:
                    # Divide em pontos se houver quebras de linha
                    if '\n' in para:
                        points = [p.strip() for p in para.split('\n') if p.strip()]
                        if current_slide:
                            current_slide['content'].extend(points)
                        else:
                            current_slide = {
                                'title': 'Conteúdo',
                                'content': points
                            }
                    else:
                        if current_slide:
                            current_slide['content'].append(para)
                        else:
                            current_slide = {
                                'title': 'Conteúdo',
                                'content': [para]
                            }
            
            # Adiciona último slide
            if current_slide:
                slides_data.append(current_slide)
            
            # Se não houver slides, cria um básico
            if not slides_data:
                slides_data.append({
                    'title': 'Conteúdo',
                    'content': [text_content]
                })
            
            return self.create_presentation(title, slides_data)
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

